"""PPTX生成器 - 全页布局，模板深度读取，预览支持"""

import os
import tempfile
from io import BytesIO
from typing import List, Optional, Tuple, Dict
from copy import deepcopy
import math

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from styles import (
    get_theme, get_style_name, SLIDE_WIDTH, SLIDE_HEIGHT,
    STYLE_CATEGORIES, LAYOUT_NAMES,
)
from parser import SlideContent
from image_search import UnsplashSearcher, translate_keywords

MARGIN = Cm(1.5)
IMAGE_FRACTION = 0.45

# EMU 单位换算: 1 inch = 914400 EMU, 1 cm = 360000 EMU
EMU_PER_CM = 360000


def create_presentation(
    slides: List[SlideContent],
    style_id: str,
    unsplash_key: str = "",
    enable_images: bool = True,
    template_prs: Optional[Presentation] = None,
) -> Presentation:
    """创建完整的PPT演示文稿"""
    theme = get_theme(style_id)

    if template_prs is not None:
        return _create_from_template(slides, theme, template_prs, style_id,
                                      unsplash_key, enable_images)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]
    searcher = UnsplashSearcher(unsplash_key) if unsplash_key else None

    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        _set_background(slide, theme["bg_color"])
        layout = slide_data.layout

        if layout == "title":
            _draw_title_slide(slide, slide_data, theme)
        elif layout == "section":
            _draw_section_slide(slide, slide_data, theme)
        elif layout == "content":
            _draw_content_slide(slide, slide_data, theme)
        elif layout == "image_text":
            _draw_content_slide(slide, slide_data, theme)
        elif layout == "toc":
            _draw_toc_slide(slide, slide_data, theme)
        elif layout == "end":
            _draw_end_slide(slide, slide_data, theme)
        else:
            _draw_content_slide(slide, slide_data, theme)

        _add_decorations(slide, theme, layout)
        if layout not in ("title", "end"):
            _add_page_number(slide, idx, len(slides), theme)

    if enable_images and searcher:
        _insert_images_builtin(prs, slides, searcher, theme)

    return prs


# ══════════════════════════════════════════════════════════════
#  模板上传生成 - 深度读取模板实际slide
# ══════════════════════════════════════════════════════════════

def _create_from_template(
    slides: List[SlideContent],
    theme: dict,
    template_prs: Presentation,
    style_id: str,
    unsplash_key: str,
    enable_images: bool,
) -> Presentation:
    """基于用户上传的模板生成PPT

    核心原则：保留模板的一切视觉元素（背景、图片、装饰、校徽等），
    只修改文本框中的文字内容。
    """
    import copy

    sw = template_prs.slide_width
    sh = template_prs.slide_height

    # ── 步骤1: 分析模板中每个slide ──
    template_slides_info = []
    for i, tslide in enumerate(template_prs.slides):
        info = _analyze_template_slide(tslide, sw, sh)
        info["index"] = i
        template_slides_info.append(info)

    # ── 步骤2: 按类型分组 ──
    cover_list = [s for s in template_slides_info if s["type"] == "cover"]
    toc_list = [s for s in template_slides_info if s["type"] == "toc"]
    section_list = [s for s in template_slides_info if s["type"] == "section"]
    content_list = [s for s in template_slides_info if s["type"] == "content"]
    end_list = [s for s in template_slides_info if s["type"] == "end"]

    # 筛选section模板：优先使用含有"PART"字样的正规章节分隔页
    part_sections = []
    for s in section_list:
        tslide = template_prs.slides[s["index"]]
        has_part = any(
            shape.has_text_frame and "part" in shape.text_frame.text.lower()
            for shape in tslide.shapes
        )
        if has_part:
            part_sections.append(s)
    if part_sections:
        section_list = part_sections

    # fallback
    if not cover_list:
        cover_list = template_slides_info[:1]
    if not section_list:
        section_list = cover_list
    if not content_list:
        content_list = template_slides_info[1:2] if len(template_slides_info) > 1 else template_slides_info[:1]
    if not end_list:
        end_list = cover_list

    # ── 步骤3: 收集所有章节名称，并清理"第X章"前缀 ──
    all_section_names = []
    section_num = 0
    for sd in slides:
        if sd.layout == "section":
            section_num += 1
            all_section_names.append(_clean_section_title(sd.title))
        elif sd.layout == "content" and sd.level <= 1:
            clean_title = _clean_section_title(sd.title)
            if not all_section_names or all_section_names[-1] != clean_title:
                section_num += 1
                all_section_names.append(clean_title)

    # ── 步骤4: 构建slide计划（自动插入目录页）──
    slide_plan = []  # [(type, data, tmpl_idx, section_num)]
    section_num = 0
    last_was_section = False
    has_toc = False

    for sd in slides:
        layout = sd.layout
        if layout == "title":
            slide_plan.append(("cover", sd, cover_list[0]["index"], 0))
        elif layout == "toc":
            slide_plan.append(("toc", sd, toc_list[0]["index"] if toc_list else content_list[0]["index"], 0))
            has_toc = True
        elif layout == "section":
            section_num += 1
            tmpl = section_list[(section_num - 1) % len(section_list)]
            slide_plan.append(("section", sd, tmpl["index"], section_num))
            last_was_section = True
        elif layout == "end":
            slide_plan.append(("end", sd, end_list[0]["index"], section_num))
        else:
            if sd.level <= 1 and not last_was_section:
                section_num += 1
                section_sd = SlideContent(layout="section", title=sd.title, level=1)
                tmpl = section_list[(section_num - 1) % len(section_list)]
                slide_plan.append(("section", section_sd, tmpl["index"], section_num))
                last_was_section = True
            slide_plan.append(("content", sd, content_list[0]["index"], section_num))
            last_was_section = False

    # 自动在封面后插入目录页（如果模板有目录页且用户未显式提供目录）
    if not has_toc and toc_list and len(slide_plan) > 0:
        toc_sd = SlideContent(layout="toc", title="目录", level=0)
        # 找到封面后的位置插入
        insert_pos = 1  # 默认在封面后
        slide_plan.insert(insert_pos, ("toc", toc_sd, toc_list[0]["index"], 0))

    # ── 步骤5: 创建新Presentation，按计划克隆模板slide ──
    new_prs = Presentation()
    new_prs.slide_width = sw
    new_prs.slide_height = sh

    for slide_type, data, tmpl_idx, sec_num in slide_plan:
        source_slide = template_prs.slides[tmpl_idx]
        new_slide = _clone_slide(new_prs, source_slide, template_prs)
        # 替换为用户内容
        if slide_type == "cover":
            _replace_cover_text(new_slide, data, theme)
        elif slide_type == "toc":
            _replace_toc_text(new_slide, data, theme, all_section_names)
        elif slide_type == "section":
            _replace_section_text(new_slide, data, theme, sec_num)
        elif slide_type == "content":
            _replace_content_text(new_slide, data, theme, all_section_names, sec_num, sw, sh)
        elif slide_type == "end":
            _replace_end_text(new_slide, data, theme)
        # 清除layout/master上的残留占位符文字
        _clear_stale_text(new_slide)

    # ── 步骤6: 添加页码（动态定位到右下角）──
    sw_cm = sw / EMU_PER_CM
    sh_cm = sh / EMU_PER_CM
    for i, (slide_type, _, _, _) in enumerate(slide_plan):
        if slide_type in ("content", "section"):
            slide = new_prs.slides[i]
            pn_left = int((sw_cm - 2.5) * EMU_PER_CM)
            pn_top = int((sh_cm - 0.8) * EMU_PER_CM)
            page_num_box = slide.shapes.add_textbox(
                pn_left, pn_top, int(Cm(2.0)), int(Cm(0.5))
            )
            p = page_num_box.text_frame.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.RIGHT

    # ── 步骤7: 插入配图（模板专用，使用slide_plan正确映射索引）──
    if enable_images and unsplash_key:
        searcher = UnsplashSearcher(unsplash_key)
        _insert_images_template(new_prs, slide_plan, searcher, theme, sw, sh)

    return new_prs


# ══════════════════════════════════════════════════════════════
#  模板文本替换函数
# ══════════════════════════════════════════════════════════════

def _clean_section_title(title: str) -> str:
    """清理章节标题：去除'第X章'前缀，只保留核心名称"""
    import re
    # 去除"第一章 ""第二章 "等中文序号前缀
    cleaned = re.sub(r'^第[一二三四五六七八九十\d]+章\s*', '', title)
    # 去除"Chapter 1 "等英文前缀
    cleaned = re.sub(r'^[Cc]hapter\s*\d+\s*', '', cleaned)
    return cleaned.strip() if cleaned.strip() else title


def _collect_text_shapes(slide):
    """收集slide中所有有文本的shape信息（位置单位为英寸，与python-pptx内部一致）"""
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue

        max_font = 0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    max_font = max(max_font, run.font.size.pt)

        # 使用英寸作为位置单位（python-pptx内部 EMU/914400 = 英寸）
        left_in = shape.left / 914400
        top_in = shape.top / 914400

        text_shapes.append({
            "shape": shape,
            "text": text,
            "font_size": max_font,
            "top": shape.top,
            "left": shape.left,
            "left_in": left_in,
            "top_in": top_in,
            "width": shape.width,
            "height": shape.height,
        })
    return text_shapes


def _replace_cover_text(slide, data: SlideContent, theme: dict):
    """封面页：主标题填入最大字号shape，副标题行智能匹配到模板对应框"""
    # 收集所有文本shape
    all_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue

        max_font = 0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    max_font = max(max_font, run.font.size.pt)

        is_ph = False
        try:
            ph = shape._element.find('.//' + qn('p:ph'))
            if ph is not None:
                is_ph = True
        except:
            pass

        all_shapes.append({
            "shape": shape,
            "text": text,
            "font_size": max_font,
            "top": shape.top,
            "left": shape.left,
            "width": shape.width,
            "height": shape.height,
            "is_placeholder": is_ph,
        })

    if not all_shapes:
        return

    # 1. 找最大字号shape填入主标题（通常是placeholder的40pt标题框）
    all_sorted = sorted(all_shapes, key=lambda x: -x["font_size"])
    _set_shape_text(all_sorted[0]["shape"], data.title)
    title_shape = all_sorted[0]["shape"]

    # 2. 副标题行智能匹配：根据模板框的文本前缀匹配用户输入
    subtitle_lines = data.subtitle.split("\n") if data.subtitle else []
    
    # 建立前缀匹配映射：模板框前缀 → 用户副标题行
    prefix_map = {
        "姓": None, "姓名": None, "姓    名": None,
        "学": None, "学号": None, "学    号": None,
        "老": None, "老师": None, "老    师": None,
    }
    
    # 为每个用户副标题行确定其类型
    typed_lines = []  # [(prefix_type, full_text)]
    for line in subtitle_lines:
        line_stripped = line.replace(" ", "")
        if line_stripped.startswith("姓") or line_stripped.startswith("姓名"):
            typed_lines.append(("name", line))
        elif line_stripped.startswith("学") or line_stripped.startswith("学号"):
            typed_lines.append(("sid", line))
        elif line_stripped.startswith("老") or line_stripped.startswith("老师"):
            typed_lines.append(("teacher", line))
        else:
            typed_lines.append(("other", line))

    # 3. 对非标题的shape，根据其原始文本前缀匹配填入
    remaining = [ts for ts in all_shapes if ts["shape"] != title_shape]
    
    for ts in remaining:
        orig_text = ts["text"].replace(" ", "")
        matched = False
        
        # 尝试匹配模板框的前缀
        if orig_text.startswith("姓") or orig_text.startswith("姓名"):
            for i, (ptype, ptext) in enumerate(typed_lines):
                if ptype == "name":
                    _set_shape_text(ts["shape"], ptext)
                    typed_lines.pop(i)
                    matched = True
                    break
        elif orig_text.startswith("学") or orig_text.startswith("学号"):
            for i, (ptype, ptext) in enumerate(typed_lines):
                if ptype == "sid":
                    _set_shape_text(ts["shape"], ptext)
                    typed_lines.pop(i)
                    matched = True
                    break
        elif orig_text.startswith("老") or orig_text.startswith("老师"):
            for i, (ptype, ptext) in enumerate(typed_lines):
                if ptype == "teacher":
                    _set_shape_text(ts["shape"], ptext)
                    typed_lines.pop(i)
                    matched = True
                    break
        
        if not matched:
            # 未匹配的框：如果有剩余副标题行则填入，否则清空
            if ts["is_placeholder"]:
                _set_shape_text(ts["shape"], "")
            elif typed_lines:
                _set_shape_text(ts["shape"], typed_lines.pop(0)[1])
            else:
                _set_shape_text(ts["shape"], "")


def _clear_stale_text(slide):
    """清除slide中来自layout/master的残留占位符文字
    （如'单击此处编辑母版标题样式'等）。
    只清除placeholder类型的shape或包含特定占位符关键词的文本。
    注意：在文本替换之后调用，避免误删已填入的用户内容。
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue

        # 检查是否包含模板占位符特征文字
        stale_keywords = [
            "单击此处编辑", "母版标题样式", "母版副标题样式",
            "母版文本样式",
        ]
        has_stale = any(kw in text for kw in stale_keywords)

        if has_stale:
            _set_shape_text(shape, "")


def _replace_toc_text(slide, data: SlideContent, theme: dict,
                       all_section_names: list):
    """目录页：保留"目录"/"CONTENTS"，大号编号更新序号，小号文本框填入章节名称
    当章节数超过模板目录位置数时，在小号框内用换行合并显示多个章节名。
    """
    text_shapes = _collect_text_shapes(slide)

    # 分三类：目录标题、大号编号(>40pt)、小号文本(≤40pt且非标题)
    title_shapes = [ts for ts in text_shapes if any(kw in ts["text"] for kw in ["目录", "CONTENTS", "目 录"])]
    number_shapes = [ts for ts in text_shapes if ts["font_size"] > 40 and ts not in title_shapes]
    small_shapes = [ts for ts in text_shapes if ts not in title_shapes and ts not in number_shapes]

    toc_items = all_section_names if all_section_names else (data.bullet_items if data.bullet_items else [])

    # 小号文本框按位置从左到右、从上到下排序
    small_sorted = sorted(small_shapes, key=lambda x: (x["left"], x["top"]))

    # 如果章节数超过小号框数量，将多个章节名合并到一个框内（用换行分隔）
    num_slots = len(small_sorted)
    if num_slots > 0 and len(toc_items) > num_slots:
        # 将toc_items均匀分配到各slot中
        per_slot = math.ceil(len(toc_items) / num_slots)
        for i, ts in enumerate(small_sorted):
            start = i * per_slot
            end = min(start + per_slot, len(toc_items))
            if start < len(toc_items):
                combined = "\n".join(toc_items[start:end])
                _set_shape_text(ts["shape"], combined)
            else:
                _set_shape_text(ts["shape"], "")
    else:
        # 正常1对1填充
        toc_idx = 0
        for ts in small_sorted:
            if toc_idx < len(toc_items):
                _set_shape_text(ts["shape"], toc_items[toc_idx])
                toc_idx += 1
            else:
                _set_shape_text(ts["shape"], "")

    # 更新大号编号框的序号
    number_sorted = sorted(number_shapes, key=lambda x: (x["left"], x["top"]))
    for i, ts in enumerate(number_sorted):
        _set_shape_text(ts["shape"], f"{i+1:02d}")


def _replace_section_text(slide, data: SlideContent, theme: dict,
                           section_num: int):
    """章节分隔页：最大→章节编号(01)，次大→'01 章节名'，PART标签根据序号更新"""
    text_shapes = _collect_text_shapes(slide)
    sorted_shapes = sorted(text_shapes, key=lambda x: -x["font_size"])

    num_str = f"{section_num:02d}"
    clean_title = _clean_section_title(data.title)
    section_title = f"{num_str} {clean_title}"

    # PART序号映射
    part_names = {1: "ONE", 2: "TWO", 3: "THREE", 4: "FOUR", 5: "FIVE",
                  6: "SIX", 7: "SEVEN", 8: "EIGHT", 9: "NINE", 10: "TEN"}
    part_name = part_names.get(section_num, str(section_num))

    for i, ts in enumerate(sorted_shapes):
        if i == 0:
            _set_shape_text(ts["shape"], num_str)
        elif i == 1:
            _set_shape_text(ts["shape"], section_title)
        else:
            if "part" in ts["text"].lower():
                _set_shape_text(ts["shape"], f"PART {part_name}")
            else:
                _set_shape_text(ts["shape"], "")


def _replace_content_text(slide, data: SlideContent, theme: dict,
                           all_section_names: list, section_num: int,
                           sw: int = 0, sh: int = 0):
    """正文内容页：侧边→章节名，标题→内容标题，动态计算正文文本框填满页面"""
    text_shapes = _collect_text_shapes(slide)
    # 侧边目录（left < 2英寸 ≈ 5cm）
    sidebar_shapes = [ts for ts in text_shapes if ts["left_in"] < 2.0]
    content_shapes = [ts for ts in text_shapes if ts["left_in"] >= 2.0]

    # 侧边目录：填入各章节名
    sidebar_sorted = sorted(sidebar_shapes, key=lambda x: x["top"])
    for i, ts in enumerate(sidebar_sorted):
        if i < len(all_section_names):
            _set_shape_text(ts["shape"], all_section_names[i])
        else:
            _set_shape_text(ts["shape"], "")

    # 正文区：最大→标题，其余清空
    content_sorted = sorted(content_shapes, key=lambda x: -x["font_size"])
    for i, ts in enumerate(content_sorted):
        if i == 0:
            _set_shape_text(ts["shape"], data.title)
        else:
            _set_shape_text(ts["shape"], "")

    # ── 动态计算正文文本框位置和大小，填满可用区域 ──
    items = data.body_lines if data.body_lines else data.bullet_items
    if not items:
        return

    # 判断是否有配图（有keywords时预留右侧空间）
    has_image = bool(data.keywords and len(data.keywords) > 0)

    if sw > 0 and sh > 0 and content_sorted:
        # 获取标题shape的实际位置来计算正文区域
        title_shape = content_sorted[0]["shape"]
        title_left = title_shape.left
        title_top = title_shape.top
        title_bottom = title_top + title_shape.height

        # 正文起始：标题下方留0.3cm间距
        body_left = title_left
        body_top = title_bottom + int(Cm(0.3))

        # 正文宽度：如果有图片则约占50%宽度，否则填满右侧
        if has_image:
            body_width = int(sw * 0.48)
        else:
            body_width = sw - body_left - int(Cm(1.0))

        # 正文高度：从标题下方到距底部1.5cm（给页码留空间）
        body_height = sh - body_top - int(Cm(1.5))

        # 安全检查：确保尺寸合理
        if body_width < int(Cm(3.0)):
            body_width = int(Cm(3.0))
        if body_height < int(Cm(2.0)):
            body_height = int(Cm(2.0))
    else:
        # 兜底默认值
        body_left = int(Cm(5.6))
        body_top = int(Cm(3.0))
        body_width = int(Cm(20.0)) if not has_image else int(Cm(14.0))
        body_height = int(Cm(13.0))

    # 根据内容量自适应字号
    total_chars = sum(len(item) for item in items)
    if total_chars > 600:
        body_font_size = Pt(14)
    elif total_chars > 350:
        body_font_size = Pt(16)
    else:
        body_font_size = Pt(18)

    txBox = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = body_font_size
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.font.name = "微软雅黑"
        p.space_after = Pt(4)
        p.space_before = Pt(2)


def _replace_end_text(slide, data: SlideContent, theme: dict):
    """结尾页：最大文本→结束语，保留其余"""
    text_shapes = _collect_text_shapes(slide)
    sorted_shapes = sorted(text_shapes, key=lambda x: -x["font_size"])
    for i, ts in enumerate(sorted_shapes):
        if i == 0:
            _set_shape_text(ts["shape"], data.title)
        else:
            pass


def _set_shape_text(shape, new_text: str):
    """替换shape中的文本，保留第一个run的格式"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        if para.runs:
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run.text = ""
        else:
            para.text = new_text


# ══════════════════════════════════════════════════════════════
#  模板分析 & 克隆
# ══════════════════════════════════════════════════════════════

def _analyze_template_slide(slide, sw: int, sh: int) -> dict:
    """分析模板中一个实际slide的结构"""
    shapes_info = []
    has_big_title = False
    has_subtitle = False
    has_body = False
    has_section_marker = False
    has_toc_marker = False
    has_end_marker = False
    title_font_size = 0
    max_font_size = 0
    has_image = False
    small_title_count = 0

    sw_in = sw / 914400
    sh_in = sh / 914400

    for shape in slide.shapes:
        if hasattr(shape, 'image'):
            has_image = True
            continue
        if not shape.has_text_frame:
            continue

        left_in = shape.left / 914400
        top_in = shape.top / 914400
        width_in = shape.width / 914400
        height_in = shape.height / 914400

        text = shape.text_frame.text.strip()
        if not text:
            continue

        shape_max_font = 0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_pt = run.font.size.pt
                    shape_max_font = max(shape_max_font, font_pt)
                    max_font_size = max(max_font_size, font_pt)

        is_top = top_in < sh_in * 0.4
        text_lower = text.lower().strip()

        if shape_max_font >= 60 and (text.isdigit() or len(text) <= 3):
            has_section_marker = True
        if "part" in text_lower:
            has_section_marker = True
        if any(kw in text_lower for kw in ["目录", "contents", "目 录"]):
            has_toc_marker = True
        if any(kw in text_lower for kw in ["谢谢", "感谢", "thank", "批评指正", "结束", "聆听"]):
            has_end_marker = True
        if shape_max_font >= 40 and is_top:
            has_big_title = True
            title_font_size = max(title_font_size, shape_max_font)
        if 16 <= shape_max_font <= 24 and not is_top:
            small_title_count += 1
        if 24 <= shape_max_font <= 32:
            has_body = True
        if any(kw in text for kw in ["姓名", "学号", "老师", "姓名："]):
            has_subtitle = True

        shapes_info.append({
            "shape": shape, "text": text,
            "left": shape.left, "top": shape.top,
            "width": shape.width, "height": shape.height,
            "font_size": shape_max_font,
            "is_top": is_top, "is_bottom": not is_top,
            "para_count": len(shape.text_frame.paragraphs),
        })

    if has_end_marker:
        slide_type = "end"
    elif has_toc_marker:
        slide_type = "toc"
    elif has_section_marker:
        slide_type = "section"
    elif has_big_title and has_image:
        slide_type = "cover"
    elif has_big_title and has_subtitle:
        slide_type = "cover"
    elif has_big_title and not has_body and not has_section_marker:
        slide_type = "cover"
    elif has_body or small_title_count >= 2:
        slide_type = "content"
    else:
        slide_type = "content"

    text_areas = []
    for si in shapes_info:
        if si["text"]:
            text_areas.append({
                "left": si["left"], "top": si["top"],
                "width": si["width"], "height": si["height"],
                "font_size": si["font_size"],
                "is_top": si["is_top"], "is_bottom": si["is_bottom"],
                "para_count": si["para_count"],
            })

    return {
        "type": slide_type, "shapes": shapes_info, "text_areas": text_areas,
        "has_big_title": has_big_title, "has_body": has_body,
        "has_image": has_image, "has_section_marker": has_section_marker,
        "small_title_count": small_title_count, "max_font_size": max_font_size,
    }


def _get_theme_colors(prs: Presentation) -> dict:
    """从模板演示文稿中提取主题颜色映射表 {scheme_name: #RRGGBB}
    包含XML中使用的别名映射(tx1→dk1, bg1→lt1等)
    """
    try:
        master = prs.slide_masters[0]
        for rel in master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_xml = etree.fromstring(rel.target_part.blob)
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                clr_scheme = theme_xml.find('.//a:clrScheme', ns)
                if clr_scheme is None:
                    return {}
                colors = {}
                for child in clr_scheme:
                    tag = child.tag.split('}')[-1]
                    for color_elem in child:
                        color_tag = color_elem.tag.split('}')[-1]
                        if color_tag == 'srgbClr':
                            colors[tag] = color_elem.get('val', '')
                        elif color_tag == 'sysClr':
                            val = color_elem.get('lastClr', '')
                            if val:
                                colors[tag] = val
                # 添加XML中常用的别名映射
                # tx1/dk1 都是深色文字, lt1/bg1 都是浅色背景
                alias_map = {
                    'tx1': colors.get('dk1', '000000'),
                    'tx2': colors.get('dk2', '44546A'),
                    'bg1': colors.get('lt1', 'FFFFFF'),
                    'bg2': colors.get('lt2', 'E7E6E6'),
                }
                colors.update(alias_map)
                return colors
    except Exception:
        pass
    return {}


def _resolve_scheme_colors(slide, theme_colors: dict):
    """将slide中所有a:schemeClr替换为a:srgbClr，防止新PPT主题色不同导致颜色错乱"""
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    # 查找slide XML中所有schemeClr元素
    for scheme_clr in slide._element.findall(f'.//{{{a_ns}}}schemeClr'):
        val = scheme_clr.get('val', '')
        if val in theme_colors:
            rgb_val = theme_colors[val]
            # 创建srgbClr替代schemeClr
            parent = scheme_clr.getparent()
            srgb_clr = etree.SubElement(parent, f'{{{a_ns}}}srgbClr')
            srgb_clr.set('val', rgb_val)
            # 保留子元素（如alpha透明度）
            for child in list(scheme_clr):
                srgb_clr.append(child)
            parent.remove(scheme_clr)


def _clone_slide(prs: Presentation, source_slide, source_prs: Presentation):
    """复制模板slide到新演示文稿（包括背景、layout形状、图片和形状）
    关键修复：
    1. 所有shape用deepcopy XML保留完整格式（裁剪/缩放/偏移），图片不变形
    2. 图片数据通过关系复制，修复跨演示文稿图片引用
    3. 将schemeClr(主题色引用)解析为srgbClr(显式RGB)
    4. 复制layout上的shape（如正文页左侧红色矩形）
    """
    import copy

    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    # ── 复制背景：先检查slide自身，再检查layout，再检查master ──
    src_cSld = source_slide._element.find(qn('p:cSld'))
    new_cSld = new_slide._element.find(qn('p:cSld'))

    bg_to_copy = None
    if src_cSld is not None:
        bg_to_copy = src_cSld.find(qn('p:bg'))
    if bg_to_copy is None:
        try:
            layout_cSld = source_slide.slide_layout._element.find(qn('p:cSld'))
            if layout_cSld is not None:
                bg_to_copy = layout_cSld.find(qn('p:bg'))
        except:
            pass
    if bg_to_copy is None:
        try:
            master_cSld = source_slide.slide_layout.slide_master._element.find(qn('p:cSld'))
            if master_cSld is not None:
                bg_to_copy = master_cSld.find(qn('p:bg'))
        except:
            pass

    if bg_to_copy is not None and new_cSld is not None:
        new_bg = new_cSld.find(qn('p:bg'))
        if new_bg is not None:
            new_cSld.remove(new_bg)
        new_cSld.insert(0, copy.deepcopy(bg_to_copy))

    # ── 复制layout和slide的shape（deepcopy保留图片裁剪/缩放）──
    _copy_shapes(source_slide.slide_layout.shapes, new_slide)
    _copy_shapes(source_slide.shapes, new_slide)

    # ── 修复图片数据：deepcopy XML后图片引用指向源文件，需重新绑定 ──
    _fix_copied_images(new_slide, source_slide, source_prs)

    # ── 解析主题色：将所有schemeClr转为srgbClr ──
    theme_colors = _get_theme_colors(source_prs)
    if theme_colors:
        _resolve_scheme_colors(new_slide, theme_colors)

    return new_slide


def _copy_shapes(shapes, new_slide):
    """复制shape列表到新slide，deepcopy XML保留完整格式（裁剪/缩放/偏移/填充）"""
    import copy
    for shape in shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)


def _fix_copied_images(new_slide, source_slide, source_prs):
    """修复deepcopy后的图片：将源演示文稿的图片数据复制到新slide的关联中
    需要同时处理来自slide自身和来自layout的图片引用。
    """
    import copy

    new_slide_part = new_slide.part
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    # 收集源slide和layout的所有图片关联
    source_image_rels = {}
    try:
        for rId, rel in source_slide.part.rels.items():
            if 'image' in rel.reltype:
                source_image_rels[rId] = rel.target_part
    except:
        pass
    try:
        for rId, rel in source_slide.slide_layout.part.rels.items():
            if 'image' in rel.reltype and rId not in source_image_rels:
                source_image_rels[rId] = rel.target_part
    except:
        pass

    # 查找新slide中所有a:blip元素（图片引用），无论其在p:blipFill还是a:blipFill中
    for blip in new_slide._element.findall(f'.//{{{a_ns}}}blip'):
        r_embed = blip.get(f'{{{r_ns}}}embed')
        if r_embed is None:
            continue

        try:
            source_image_part = source_image_rels.get(r_embed)
            if source_image_part is None:
                continue
            img_blob = source_image_part.blob
            img_stream = BytesIO(img_blob)
            new_image_part, new_rId = new_slide_part.get_or_add_image_part(img_stream)
            blip.set(f'{{{r_ns}}}embed', new_rId)
        except Exception:
            pass


# ══════════════════════════════════════════════════════════════
#  图片插入
# ══════════════════════════════════════════════════════════════

def _insert_images_template(prs, slide_plan, searcher, theme, sw, sh):
    """模板模式的图片插入——使用slide_plan正确映射索引，图片占右侧大面积"""
    for plan_idx, (slide_type, data, tmpl_idx, sec_num) in enumerate(slide_plan):
        if slide_type != "content":
            continue
        if not data.keywords:
            continue

        query = translate_keywords(data.keywords) or data.image_query
        if not query:
            continue
        results = searcher.search_images(query, per_page=1, orientation="landscape")
        if not results:
            continue
        img_data = searcher.download_image(results[0]["url"])
        if not img_data:
            continue

        try:
            slide = prs.slides[plan_idx]
            img_stream = BytesIO(img_data)

            # 图片放在右侧，占约45%页面宽度、75%页面高度
            img_left = int(sw * 0.53)
            img_top = int(sh * 0.12)
            img_width = int(sw * 0.43)
            img_height = int(sh * 0.72)

            slide.shapes.add_picture(img_stream, img_left, img_top, img_width, img_height)
        except Exception as e:
            print(f"插入图片失败: {e}")


def _insert_images_builtin(prs, slides, searcher, theme):
    """内置风格的图片插入——图片占右侧大面积"""
    for i, slide_data in enumerate(slides):
        if slide_data.layout not in ("image_text", "content"):
            continue
        if not slide_data.image_query:
            continue
        query = translate_keywords(slide_data.keywords) or slide_data.image_query
        results = searcher.search_images(query, per_page=1, orientation="landscape")
        if not results:
            continue
        img_data = searcher.download_image(results[0]["url"])
        if not img_data:
            continue
        try:
            img_stream = BytesIO(img_data)
            slide = prs.slides[i]
            # 图片占右侧约40%宽度、70%高度
            img_left = int(prs.slide_width * 0.57)
            img_top = int(prs.slide_height * 0.15)
            img_width = int(prs.slide_width * 0.38)
            img_height = int(prs.slide_height * 0.68)
            slide.shapes.add_picture(img_stream, img_left, img_top, img_width, img_height)
        except Exception as e:
            print(f"插入图片失败: {e}")


# ══════════════════════════════════════════════════════════════
#  预览功能 - 生成幻灯片缩略图
# ══════════════════════════════════════════════════════════════

def generate_preview_html(prs: Presentation) -> str:
    """生成PPT预览的HTML，直接在浏览器中渲染（解决中文乱码问题）"""
    slides_html = []
    sw_cm = prs.slide_width / EMU_PER_CM
    sh_cm = prs.slide_height / EMU_PER_CM

    for page_idx, slide in enumerate(prs.slides):
        bg_color = "#FFFFFF"
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                bg_rgb = fill.fore_color.rgb
                bg_color = f"#{str(bg_rgb)[0:2]}{str(bg_rgb)[2:4]}{str(bg_rgb)[4:6]}"
        except:
            pass

        shapes_html = []
        for shape in slide.shapes:
            left_pct = shape.left / prs.slide_width * 100
            top_pct = shape.top / prs.slide_height * 100
            width_pct = shape.width / prs.slide_width * 100
            height_pct = shape.height / prs.slide_height * 100

            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                font_size = 14
                font_color = "#333333"
                font_bold = False
                try:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                font_size = max(font_size, int(run.font.size.pt * 0.55))
                            if run.font.color and run.font.color.rgb:
                                c = run.font.color.rgb
                                font_color = f"#{str(c)[0:2]}{str(c)[2:4]}{str(c)[4:6]}"
                            if run.font.bold:
                                font_bold = True
                except:
                    pass

                text_escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                lines = text_escaped.split('\n')
                display_lines = '<br>'.join(lines[:8])
                if len(lines) > 8:
                    display_lines += '...'

                bold_css = "font-weight:bold;" if font_bold else ""
                shapes_html.append(
                    f'<div style="position:absolute;left:{left_pct}%;top:{top_pct}%;'
                    f'width:{width_pct}%;height:{height_pct}%;'
                    f'font-size:{font_size}px;color:{font_color};{bold_css}'
                    f'overflow:hidden;word-wrap:break-word;padding:2px;">'
                    f'{display_lines}</div>'
                )

            elif hasattr(shape, 'image'):
                shapes_html.append(
                    f'<div style="position:absolute;left:{left_pct}%;top:{top_pct}%;'
                    f'width:{width_pct}%;height:{height_pct}%;'
                    f'background:#e0e0e0;border-radius:4px;display:flex;'
                    f'align-items:center;justify-content:center;color:#999;font-size:10px;">'
                    f'图片</div>'
                )

        slides_html.append(
            f'<div style="position:relative;width:100%;padding-bottom:{sh_cm/sw_cm*100}%;'
            f'background:{bg_color};border:1px solid #ddd;border-radius:4px;margin-bottom:8px;">'
            f'{"".join(shapes_html)}</div>'
        )

    return "".join(slides_html)


def generate_preview_images(prs: Presentation) -> List[bytes]:
    """生成PPT每页的预览图 - 已废弃，改用HTML预览"""
    return []


# ══════════════════════════════════════════════════════════════
#  内置风格绘制 — 全页填充布局 + 图文混排
# ══════════════════════════════════════════════════════════════

def _draw_title_slide(slide, data, theme):
    dc = theme["decorative_colors"]
    _add_shape(slide, 0, SLIDE_HEIGHT - Cm(5.0), SLIDE_WIDTH, Cm(5.0), theme["bg_color2"], alpha=0.6)
    _add_shape(slide, 0, 0, SLIDE_WIDTH, Cm(0.15), dc[0])

    txBox = _add_textbox(slide, Cm(2.0), Cm(2.0), SLIDE_WIDTH - Cm(4.0), Cm(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    p.font.name = theme["title_font"]
    p.alignment = PP_ALIGN.CENTER

    _add_shape(slide, SLIDE_WIDTH // 2 - Cm(3.0), Cm(6.8), Cm(6.0), Cm(0.06), theme["accent_color"])
    _add_shape_diamond(slide, SLIDE_WIDTH // 2 - Cm(0.2), Cm(6.65), Cm(0.4), Cm(0.4), theme["accent_color"])

    if data.subtitle:
        txBox2 = _add_textbox(slide, Cm(2.0), Cm(7.3), SLIDE_WIDTH - Cm(4.0), Cm(2.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = data.subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = theme["accent_color"]
        p2.font.name = theme["body_font"]
        p2.alignment = PP_ALIGN.CENTER


def _draw_section_slide(slide, data, theme):
    dc = theme["decorative_colors"]
    _add_shape(slide, 0, 0, Cm(1.0), SLIDE_HEIGHT, theme["accent_color"])
    _add_shape(slide, Cm(1.0), 0, Cm(0.08), SLIDE_HEIGHT, dc[1], alpha=0.5)
    _add_shape(slide, Cm(5.0), Cm(1.0), SLIDE_WIDTH - Cm(5.5), Cm(9.0), theme["bg_color2"], alpha=0.4)

    txBox_label = _add_textbox(slide, Cm(2.5), Cm(1.5), Cm(3.0), Cm(1.0))
    p_label = txBox_label.text_frame.paragraphs[0]
    p_label.text = "CHAPTER"
    p_label.font.size = Pt(12)
    p_label.font.color.rgb = theme["accent2_color"]
    p_label.font.name = theme["body_font"]
    p_label.font.bold = True

    txBox = _add_textbox(slide, Cm(2.5), Cm(2.8), SLIDE_WIDTH - Cm(4.0), Cm(6.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.title
    p.font.size = theme["section_size"]
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    p.font.name = theme["title_font"]

    _add_shape(slide, Cm(2.5), Cm(8.5), Cm(5.0), Cm(0.05), theme["accent2_color"])
    _add_shape(slide, Cm(7.5), Cm(8.35), Cm(0.3), Cm(0.3), theme["accent_color"])


def _draw_content_slide(slide, data, theme):
    """内容页：有配图时左侧文字+右侧图片区域，无配图时文字满宽"""
    dc = theme["decorative_colors"]
    _add_shape(slide, 0, 0, SLIDE_WIDTH, Cm(0.35), theme["accent_color"])
    _add_shape(slide, 0, Cm(0.35), SLIDE_WIDTH, Cm(0.06), dc[3] if len(dc)>3 else dc[0])

    # 判断是否有配图（预留右侧空间）
    has_image = bool(data.keywords and len(data.keywords) > 0)

    # 标题
    txBox = _add_textbox(slide, Cm(1.5), Cm(0.7), SLIDE_WIDTH - Cm(3.0), Cm(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    p.font.name = theme["title_font"]

    _add_shape(slide, Cm(1.0), Cm(0.7), Cm(0.12), Cm(1.5), theme["accent_color"])
    _add_shape(slide, Cm(1.5), Cm(2.6), Cm(4.0), Cm(0.04), theme["accent_color"])

    body_top = Cm(3.3)
    body_height = SLIDE_HEIGHT - Cm(4.5)

    if has_image:
        # 图文混排：文字占左55%，图片占右40%
        text_width = SLIDE_WIDTH * 52 // 100
        _add_shape(slide, Cm(1.2), Cm(2.9), text_width - Cm(1.2), body_height + Cm(0.4),
                   theme["bg_color2"], alpha=0.35, corner_radius=Cm(0.3))
        # 右侧图片区域背景
        img_area_left = SLIDE_WIDTH * 56 // 100
        img_area_width = SLIDE_WIDTH - img_area_left - Cm(0.8)
        _add_shape(slide, img_area_left, Cm(2.9), img_area_width, body_height + Cm(0.4),
                   theme["bg_color2"], alpha=0.15, corner_radius=Cm(0.3))
        actual_text_width = text_width - Cm(1.5)
    else:
        # 文字满宽
        _add_shape(slide, Cm(1.2), Cm(2.9), SLIDE_WIDTH - Cm(2.4), body_height + Cm(0.4),
                   theme["bg_color2"], alpha=0.35, corner_radius=Cm(0.3))
        actual_text_width = SLIDE_WIDTH - Cm(4.0)

    if data.bullet_items:
        txBox2 = _add_textbox(slide, Cm(2.0), body_top, actual_text_width, body_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, item in enumerate(data.bullet_items):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"  •  {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = theme["body_color"]
            p.font.name = theme["body_font"]
            p.space_after = Pt(12)
            p.space_before = Pt(4)

    elif data.body_lines:
        txBox2 = _add_textbox(slide, Cm(2.0), body_top, actual_text_width, body_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(data.body_lines):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = theme["body_color"]
            p.font.name = theme["body_font"]
            p.space_after = Pt(10)
            p.space_before = Pt(4)


def _draw_toc_slide(slide, data, theme):
    txBox = _add_textbox(slide, Cm(1.5), Cm(1.0), SLIDE_WIDTH - Cm(3.0), Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data.title or "目录"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    p.font.name = theme["title_font"]

    _add_shape(slide, Cm(1.5), Cm(2.5), Cm(3.0), Cm(0.04), theme["accent_color"])

    txBox2 = _add_textbox(slide, Cm(2.0), Cm(3.0), SLIDE_WIDTH - Cm(4.0), Cm(12.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(data.bullet_items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        run_num = p.add_run()
        run_num.text = f"  {i+1}.  "
        run_num.font.size = Pt(24)
        run_num.font.color.rgb = theme["accent_color"]
        run_num.font.name = theme["body_font"]
        run_num.font.bold = True
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(22)
        run_text.font.color.rgb = theme["body_color"]
        run_text.font.name = theme["body_font"]
        p.space_after = Pt(16)


def _draw_end_slide(slide, data, theme):
    dc = theme["decorative_colors"]
    _add_shape(slide, 0, SLIDE_HEIGHT - Cm(5.0), SLIDE_WIDTH, Cm(5.0), theme["bg_color2"], alpha=0.5)
    _add_shape(slide, Cm(2.0), Cm(5.8), SLIDE_WIDTH - Cm(4.0), Cm(0.06), theme["accent_color"])
    _add_shape_diamond(slide, SLIDE_WIDTH // 2 - Cm(0.3), Cm(5.65), Cm(0.6), Cm(0.6), theme["accent_color"])

    txBox = _add_textbox(slide, Cm(2.0), Cm(2.5), SLIDE_WIDTH - Cm(4.0), Cm(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]
    p.font.name = theme["title_font"]
    p.alignment = PP_ALIGN.CENTER

    txBox2 = _add_textbox(slide, Cm(2.0), Cm(6.8), SLIDE_WIDTH - Cm(4.0), Cm(2.0))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = "THANK YOU"
    p2.font.size = Pt(20)
    p2.font.color.rgb = theme["accent_color"]
    p2.font.name = theme["body_font"]
    p2.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════
#  装饰系统
# ══════════════════════════════════════════════════════════════

def _add_decorations(slide, theme, layout):
    ds = theme["deco_style"]
    dc = theme["decorative_colors"]
    deco_map = {
        "ink": _deco_ink, "organic": _deco_organic, "chinese": _deco_chinese,
        "bold": _deco_bold, "wood": _deco_wood, "corporate": _deco_corporate,
        "sakura": _deco_sakura, "neon": _deco_neon,
    }
    fn = deco_map.get(ds)
    if fn:
        fn(slide, dc, theme, layout)

def _deco_ink(slide, dc, theme, layout):
    if layout not in ("title", "end"):
        _add_shape(slide, 0, SLIDE_HEIGHT - Cm(0.6), SLIDE_WIDTH, Cm(0.6), dc[0], alpha=0.25)
    _add_shape(slide, SLIDE_WIDTH - Cm(2.5), 0, Cm(2.5), Cm(0.6), dc[2], alpha=0.12)
    _add_shape_oval(slide, SLIDE_WIDTH - Cm(3.0), SLIDE_HEIGHT - Cm(2.5), Cm(2.5), Cm(2.0), dc[0], alpha=0.1)

def _deco_organic(slide, dc, theme, layout):
    if layout not in ("title", "end"):
        _add_shape(slide, 0, SLIDE_HEIGHT - Cm(0.5), SLIDE_WIDTH, Cm(0.5), dc[0], alpha=0.4)
    _add_shape(slide, 0, Cm(1.5), Cm(0.25), Cm(4.0), dc[1], alpha=0.3)
    _add_shape_oval(slide, SLIDE_WIDTH - Cm(4.0), -Cm(2.0), Cm(5.0), Cm(5.0), dc[3] if len(dc)>3 else dc[0], alpha=0.1)

def _deco_chinese(slide, dc, theme, layout):
    if layout not in ("title",):
        _add_shape(slide, 0, SLIDE_HEIGHT - Cm(0.25), SLIDE_WIDTH, Cm(0.25), dc[0], alpha=0.7)
    _add_shape(slide, SLIDE_WIDTH - Cm(1.0), 0, Cm(1.0), Cm(1.0), dc[1], alpha=0.2)
    _add_shape(slide, 0, SLIDE_HEIGHT - Cm(1.0), Cm(1.0), Cm(1.0), dc[1], alpha=0.15)

def _deco_bold(slide, dc, theme, layout):
    if layout not in ("title", "end"):
        _add_shape(slide, 0, SLIDE_HEIGHT - Cm(1.0), SLIDE_WIDTH, Cm(1.0), dc[0], alpha=0.85)
    _add_shape(slide, 0, 0, Cm(0.35), SLIDE_HEIGHT, dc[1], alpha=0.6)
    _add_shape(slide, SLIDE_WIDTH - Cm(2.0), 0, Cm(2.0), Cm(1.5), dc[0], alpha=0.2)

def _deco_wood(slide, dc, theme, layout):
    if layout not in ("title", "end"):
        _add_shape(slide, 0, SLIDE_HEIGHT - Cm(0.5), SLIDE_WIDTH, Cm(0.5), dc[0], alpha=0.5)
    _add_shape(slide, SLIDE_WIDTH - Cm(0.12), 0, Cm(0.12), SLIDE_HEIGHT, dc[1], alpha=0.25)
    _add_shape_oval(slide, SLIDE_WIDTH - Cm(2.0), SLIDE_HEIGHT - Cm(2.0), Cm(1.5), Cm(1.5), dc[2], alpha=0.12)

def _deco_corporate(slide, dc, theme, layout):
    if layout not in ("title", "end"):
        _add_shape(slide, 0, SLIDE_HEIGHT - Cm(0.4), SLIDE_WIDTH, Cm(0.4), dc[1], alpha=0.6)
    _add_shape(slide, 0, 0, Cm(0.8), Cm(0.8), dc[1], alpha=0.5)

def _deco_sakura(slide, dc, theme, layout):
    if layout not in ("title", "end"):
        _add_shape(slide, 0, SLIDE_HEIGHT - Cm(0.4), SLIDE_WIDTH, Cm(0.4), dc[0], alpha=0.4)
    _add_shape_oval(slide, SLIDE_WIDTH - Cm(3.5), -Cm(1.0), Cm(4.0), Cm(4.0), dc[0], alpha=0.12)
    _add_shape_oval(slide, Cm(20.0), Cm(3.0), Cm(0.4), Cm(0.4), dc[1], alpha=0.1)

def _deco_neon(slide, dc, theme, layout):
    _add_shape(slide, 0, SLIDE_HEIGHT - Cm(0.3), SLIDE_WIDTH, Cm(0.3), dc[0], alpha=0.7)
    _add_shape(slide, 0, 0, SLIDE_WIDTH, Cm(0.15), dc[1], alpha=0.5)
    _add_shape(slide, 0, 0, Cm(0.15), SLIDE_HEIGHT, dc[0], alpha=0.4)


# ══════════════════════════════════════════════════════════════
#  辅助函数
# ══════════════════════════════════════════════════════════════

def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(int(left), int(top), int(width), int(height))

def _add_shape(slide, left, top, width, height, fill_color, alpha=1.0, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        int(left), int(top), int(width), int(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha < 1.0:
        _set_shape_alpha(shape, alpha)
    return shape

def _add_shape_oval(slide, left, top, width, height, fill_color, alpha=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(left), int(top), int(width), int(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha < 1.0:
        _set_shape_alpha(shape, alpha)
    return shape

def _add_shape_diamond(slide, left, top, width, height, fill_color, alpha=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, int(left), int(top), int(width), int(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha < 1.0:
        _set_shape_alpha(shape, alpha)
    return shape

def _set_shape_alpha(shape, alpha):
    try:
        fill = shape.fill._fill
        srgb = fill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_elem = srgb.makeelement(qn('a:alpha'), {})
            alpha_elem.set('val', str(int(alpha * 100000)))
            srgb.append(alpha_elem)
    except:
        pass

def _set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_page_number(slide, idx, total, theme):
    txBox = _add_textbox(slide, SLIDE_WIDTH - Cm(2.5), SLIDE_HEIGHT - Cm(0.8), Cm(2.0), Cm(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{idx + 1} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = theme["accent_color"]
    p.font.name = theme["body_font"]
    p.alignment = PP_ALIGN.RIGHT

def save_presentation(prs, filepath):
    prs.save(filepath)

def export_to_bytes(prs):
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
